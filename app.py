# app.py
"""
SciSports Scouting Form Generator

- Search: input + Search (Enter submits), then select player. No limit/offset controls.
- Season stats: FIXED season label TARGET_SEASON_LABEL ("2025/2026") and aggregated like SciSports UI:
  - fetch all seasonIds matching that label
  - query career-stats for those seasonIds
  - if a "Total" competition row exists, use it
  - else: for each competition take MAX (cumulative rows), then sum across competitions
- Career totals: apply the same per-season logic across all seasons, then sum seasons.

Fills placeholders including {MV}.
"""

from __future__ import annotations

import json
import os
import re
import tempfile
from dataclasses import dataclass
from datetime import datetime
from typing import Any, Dict, List, Optional, Tuple

import requests
import streamlit as st
from pptx import Presentation
from pptx.dml.color import RGBColor

API_BASE = "https://api-recruitment.scisports.app/api"
TOKEN_URL = "https://identity.scisports.app/connect/token"
TEMPLATE_PATH = "TemplateScoutingsRapport.pptx"

SEARCH_LIMIT = 50
TARGET_SEASON_LABEL = "2025/2026"

SEASON_RE = re.compile(r"\b(20\d{2})\s*[/\-]\s*(\d{2}|20\d{2})\b")


@dataclass(frozen=True)
class PlayerOption:
    player_id: int
    name: str
    age: Optional[int]
    position: str
    club: str
    league: str

    def label(self) -> str:
        age_str = "?" if self.age is None else str(self.age)
        pos = self.position or "UnknownPos"
        club = self.club or "UnknownClub"
        league = self.league or "UnknownLeague"
        return f"{self.name} — {age_str} — {pos} — {club} ({league})"


# -----------------------------
# Helpers
# -----------------------------
def normalize_season_label(name: str) -> str:
    if not name:
        return ""
    m = SEASON_RE.search(str(name))
    if not m:
        return ""
    y1 = int(m.group(1))
    y2_raw = m.group(2)
    if len(y2_raw) == 2:
        y2 = (y1 // 100) * 100 + int(y2_raw)
        if y2 < y1:
            y2 += 100
    else:
        y2 = int(y2_raw)
    return f"{y1}/{y2}"


def _http_session() -> requests.Session:
    s = requests.Session()
    s.headers.update({"Accept": "application/json"})
    return s


def _auth_headers(token: str) -> Dict[str, str]:
    return {"Authorization": f"Bearer {token}"}


def _as_text(v: Any) -> str:
    return "" if v is None else str(v)


def _fmt_int(v: Any) -> str:
    try:
        return "" if v is None else str(int(v))
    except Exception:
        return ""


def _parse_iso_date_to_ddmmyyyy(value: Optional[str]) -> str:
    if not value:
        return ""
    try:
        dt = datetime.fromisoformat(value.replace("Z", "+00:00"))
        return dt.strftime("%d-%m-%Y")
    except Exception:
        try:
            dt = datetime.strptime(value[:10], "%Y-%m-%d")
            return dt.strftime("%d-%m-%Y")
        except Exception:
            return value


def _fmt_money_eur(value: Any) -> str:
    try:
        if value is None:
            return ""
        v = float(value)
        if abs(v) >= 1_000_000:
            return f"€ {v/1_000_000:.2f}M"
        if abs(v) >= 1_000:
            return f"€ {v/1_000:.0f}K"
        return f"€ {v:.0f}"
    except Exception:
        return ""


def _first_position(info: Dict[str, Any]) -> str:
    positions = info.get("positions") or []
    if isinstance(positions, list) and positions:
        return _as_text(positions[0])
    return ""


def _extract_int(d: Dict[str, Any], *paths: str) -> int:
    for p in paths:
        cur: Any = d
        ok = True
        for part in p.split("."):
            if not isinstance(cur, dict) or part not in cur:
                ok = False
                break
            cur = cur[part]
        if ok and cur is not None:
            try:
                return int(round(float(cur)))
            except Exception:
                pass
    return 0


# -----------------------------
# Secrets + token
# -----------------------------
def _require_secrets() -> Dict[str, str]:
    required = [
        "SCISPORTS_USERNAME",
        "SCISPORTS_PASSWORD",
        "SCISPORTS_CLIENT_ID",
        "SCISPORTS_CLIENT_SECRET",
    ]
    missing = [k for k in required if not st.secrets.get(k)]
    if missing:
        st.error(
            "Missing Streamlit Secrets. Add these keys in Streamlit Cloud → App → Settings → Secrets:\n"
            + "\n".join([f"- {k}" for k in missing])
        )
        st.stop()

    return {
        "username": st.secrets["SCISPORTS_USERNAME"],
        "password": st.secrets["SCISPORTS_PASSWORD"],
        "client_id": st.secrets["SCISPORTS_CLIENT_ID"],
        "client_secret": st.secrets["SCISPORTS_CLIENT_SECRET"],
        "scope": st.secrets.get("SCISPORTS_SCOPE", "api recruitment"),
    }


def token_password_grant_from_secrets(timeout_s: int = 30) -> str:
    creds = _require_secrets()
    payload = {
        "grant_type": "password",
        "username": creds["username"],
        "password": creds["password"],
        "client_id": creds["client_id"],
        "client_secret": creds["client_secret"],
        "scope": creds["scope"],
    }
    resp = requests.post(TOKEN_URL, data=payload, timeout=timeout_s)
    resp.raise_for_status()
    data = resp.json()
    token = data.get("access_token")
    if not token:
        raise RuntimeError(f"Token response missing access_token: {json.dumps(data)[:500]}")
    return token


# -----------------------------
# Pagination helper
# -----------------------------
def fetch_all_items(
    token: str,
    path: str,
    params: Dict[str, Any],
    *,
    page_limit: int = 200,
    hard_cap: int = 50_000,
) -> List[Dict[str, Any]]:
    s = _http_session()
    out: List[Dict[str, Any]] = []
    offset = int(params.get("offset", 0))
    limit = min(max(int(params.get("limit", page_limit)), 1), page_limit)

    while True:
        p = dict(params)
        p["offset"] = offset
        p["limit"] = limit

        url = f"{API_BASE}{path}"
        resp = s.get(url, headers=_auth_headers(token), params=p, timeout=30)
        resp.raise_for_status()
        payload = resp.json()

        items = payload.get("items") or []
        if not isinstance(items, list):
            break

        out.extend([it for it in items if isinstance(it, dict)])

        total = payload.get("total")
        if isinstance(total, int) and len(out) >= total:
            break
        if not items:
            break

        offset += limit
        if len(out) >= hard_cap:
            break

    return out


# -----------------------------
# API calls
# -----------------------------
@st.cache_data(show_spinner=False, ttl=60 * 15)
def search_players(token: str, search_text: str) -> Tuple[int, List[PlayerOption]]:
    s = _http_session()
    params: Dict[str, Any] = {"offset": 0, "limit": SEARCH_LIMIT}
    if search_text.strip():
        params["searchText"] = search_text.strip()

    url = f"{API_BASE}/v2/players"
    resp = s.get(url, headers=_auth_headers(token), params=params, timeout=30)
    resp.raise_for_status()
    payload = resp.json()

    total = int(payload.get("total", 0))
    items = payload.get("items") or []

    options: List[PlayerOption] = []
    for it in items:
        info = it.get("info") or {}
        team = it.get("team") or {}
        league = it.get("league") or {}
        pid = info.get("id")
        if pid is None:
            continue
        options.append(
            PlayerOption(
                player_id=int(pid),
                name=_as_text(info.get("name") or info.get("footballName") or ""),
                age=info.get("age"),
                position=_first_position(info),
                club=_as_text(team.get("name") or ""),
                league=_as_text(league.get("name") or ""),
            )
        )

    return total, options


def get_player(token: str, player_id: int) -> Dict[str, Any]:
    s = _http_session()
    url = f"{API_BASE}/v2/players/{player_id}"
    resp = s.get(url, headers=_auth_headers(token), timeout=30)
    resp.raise_for_status()
    return resp.json()


def get_latest_transfer_fee(token: str, player_id: int) -> Optional[Dict[str, Any]]:
    s = _http_session()
    url = f"{API_BASE}/v2/metrics/players/transfer-fees"
    params: Dict[str, Any] = {"offset": 0, "limit": 1, "playerIds": player_id, "latestTransferFee": "true"}
    resp = s.get(url, headers=_auth_headers(token), params=params, timeout=30)
    resp.raise_for_status()
    data = resp.json()
    items = data.get("items") or []
    return items[0] if items else None


@st.cache_data(show_spinner=False, ttl=60 * 60)
def get_seasons_for_player(token: str, player_id: int) -> List[Dict[str, Any]]:
    # Some SciSports endpoints behave differently with casing -> send both.
    return fetch_all_items(
        token,
        "/v2/seasons",
        params={"offset": 0, "limit": 200, "playerIds": player_id, "PlayerIds": player_id},
        page_limit=200,
        hard_cap=5000,
    )


def season_ids_for_label(seasons: List[Dict[str, Any]], label: str) -> List[int]:
    target = normalize_season_label(label)
    out: List[int] = []
    for it in seasons:
        sid = it.get("id")
        if not isinstance(sid, int):
            continue
        nm = normalize_season_label(_as_text(it.get("name") or ""))
        if nm == target:
            out.append(sid)
    return sorted(list(set(out)))


def _competition_name_lower(it: Dict[str, Any]) -> str:
    comp = it.get("competition") or it.get("competitionGroup") or it.get("league") or {}
    if isinstance(comp, dict):
        name = _as_text(comp.get("name") or "").strip().lower()
        if name:
            return name
    return _as_text(it.get("competitionName") or it.get("leagueName") or "").strip().lower()


def _competition_key(it: Dict[str, Any]) -> str:
    comp = it.get("competition") or it.get("competitionGroup") or it.get("league") or {}
    if isinstance(comp, dict):
        cid = comp.get("id")
        if cid is not None:
            return f"id:{cid}"
        name = _as_text(comp.get("name") or "").strip().lower()
        if name:
            return f"name:{name}"
    name2 = _as_text(it.get("competitionName") or it.get("leagueName") or "").strip().lower()
    return f"name:{name2}" if name2 else "unknown"


def _is_total_row(it: Dict[str, Any]) -> bool:
    name = _competition_name_lower(it)
    return name in {
        "total",
        "overall",
        "all",
        "all competitions",
        "all competitions total",
    }


def aggregate_career_stats_items_like_ui(items: List[Dict[str, Any]]) -> Dict[str, int]:
    """
    Mimics the safe aggregation your old script used:
    - Prefer a Total row if present (highest minutes wins)
    - Else: MAX per competition (cumulative rows), then SUM competitions
    """
    if not items:
        return {"matches": 0, "minutes": 0, "goals": 0, "assists": 0}

    total_rows = [it for it in items if _is_total_row(it)]
    if total_rows:
        best = max(
            total_rows,
            key=lambda it: _extract_int(it, "stats.minutesPlayed", "stats.minutes"),
        )
        return {
            "matches": _extract_int(best, "stats.matchesPlayed", "stats.matches", "stats.games"),
            "minutes": _extract_int(best, "stats.minutesPlayed", "stats.minutes"),
            "goals": _extract_int(best, "stats.goal", "stats.goals", "stats.goalNonPenalty"),
            "assists": _extract_int(best, "stats.assist", "stats.assists"),
        }

    best_per_comp: Dict[str, Dict[str, int]] = {}
    for it in items:
        ck = _competition_key(it)
        row = {
            "matches": _extract_int(it, "stats.matchesPlayed", "stats.matches", "stats.games"),
            "minutes": _extract_int(it, "stats.minutesPlayed", "stats.minutes"),
            "goals": _extract_int(it, "stats.goal", "stats.goals", "stats.goalNonPenalty"),
            "assists": _extract_int(it, "stats.assist", "stats.assists"),
        }
        prev = best_per_comp.get(ck)
        if prev is None:
            best_per_comp[ck] = row
        else:
            prev["matches"] = max(prev["matches"], row["matches"])
            prev["minutes"] = max(prev["minutes"], row["minutes"])
            prev["goals"] = max(prev["goals"], row["goals"])
            prev["assists"] = max(prev["assists"], row["assists"])

    return {
        "matches": sum(v["matches"] for v in best_per_comp.values()),
        "minutes": sum(v["minutes"] for v in best_per_comp.values()),
        "goals": sum(v["goals"] for v in best_per_comp.values()),
        "assists": sum(v["assists"] for v in best_per_comp.values()),
    }


def get_career_stats_for_season_ids(token: str, player_id: int, season_ids: List[int]) -> List[Dict[str, Any]]:
    if not season_ids:
        return []
    return fetch_all_items(
        token,
        "/v2/metrics/career-stats/players",
        params={"offset": 0, "limit": 200, "playerIds": player_id, "seasonIds": season_ids},
        page_limit=200,
        hard_cap=20_000,
    )


def compute_target_season_stats(token: str, player_id: int, season_label: str) -> Tuple[str, List[int], Dict[str, int]]:
    seasons = get_seasons_for_player(token, player_id)
    season_ids = season_ids_for_label(seasons, season_label)
    items = get_career_stats_for_season_ids(token, player_id, season_ids)
    return season_label, season_ids, aggregate_career_stats_items_like_ui(items)


def compute_career_totals(token: str, player_id: int) -> Dict[str, int]:
    items = fetch_all_items(
        token,
        "/v2/metrics/career-stats/players",
        params={"offset": 0, "limit": 200, "playerIds": player_id},
        page_limit=200,
        hard_cap=50_000,
    )

    by_season: Dict[int, List[Dict[str, Any]]] = {}
    for it in items:
        sid = it.get("seasonId") or (it.get("season") or {}).get("id")
        if isinstance(sid, int):
            by_season.setdefault(sid, []).append(it)

    totals = {"matches": 0, "minutes": 0, "goals": 0, "assists": 0}
    for _sid, season_items in by_season.items():
        s = aggregate_career_stats_items_like_ui(season_items)
        totals["matches"] += s["matches"]
        totals["minutes"] += s["minutes"]
        totals["goals"] += s["goals"]
        totals["assists"] += s["assists"]
    return totals


# -----------------------------
# Agency/Agent best-effort extraction
# -----------------------------
def extract_agent_and_agency(player_obj: Dict[str, Any]) -> Tuple[str, str]:
    contract = player_obj.get("contract") or {}
    info = player_obj.get("info") or {}
    agency = (
        _as_text(contract.get("agencyName"))
        or _as_text((contract.get("agency") or {}).get("name") if isinstance(contract.get("agency"), dict) else "")
        or _as_text(info.get("agencyName"))
    )
    agent = (
        _as_text(contract.get("agentName"))
        or _as_text((contract.get("agent") or {}).get("name") if isinstance(contract.get("agent"), dict) else "")
        or _as_text(info.get("agentName"))
    )
    return agency.strip(), agent.strip()


# -----------------------------
# PPTX replacement (split-runs safe)
# -----------------------------
def _replace_tokens_in_shape(shape, replacements: Dict[str, str]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return

    for paragraph in shape.text_frame.paragraphs:
        if not paragraph.runs:
            continue

        for run in paragraph.runs:
            t = run.text or ""
            if "{" not in t:
                continue
            new_t = t
            for k, v in replacements.items():
                if k in new_t:
                    new_t = new_t.replace(k, v)
            if new_t != t:
                run.text = new_t

        combined = "".join((r.text or "") for r in paragraph.runs)
        if "{" not in combined:
            continue

        new_combined = combined
        for k, v in replacements.items():
            if k in new_combined:
                new_combined = new_combined.replace(k, v)

        if new_combined != combined:
            paragraph.runs[0].text = new_combined
            for r in paragraph.runs[1:]:
                r.text = ""


def fill_pptx(template_path: str, output_path: str, replacements: Dict[str, str], positions: List[str], preferred_foot: str) -> None:
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            _replace_tokens_in_shape(shape, replacements)
        apply_position_coloring(prs, slide, positions, preferred_foot)
    prs.save(output_path)


# -----------------------------
# Position coloring (bottom-left 1..11)
# -----------------------------
MAIN_BLUE = RGBColor(0, 83, 159)
SECOND_BLUE = RGBColor(0, 142, 204)

POSITION_TO_NUMBER: Dict[str, int] = {
    "Goalkeeper": 1,
    "RightBack": 2,
    "RightFullback": 2,
    "RightWing": 7,
    "LeftBack": 5,
    "LeftWing": 11,
    "CentreMidfield": 8,
    "AttackingMidfield": 10,
    "DefensiveMidfield": 6,
    "CentreForward": 9,
    "CentreBack": -1,
    "Centre Back": -1,
}


def _resolve_position_number(position: str, preferred_foot: str) -> Optional[int]:
    if position in {"CentreBack", "Centre Back"}:
        pf = (preferred_foot or "").strip().lower()
        is_left = "left" in pf or pf.startswith("l")
        return 4 if is_left else 3
    n = POSITION_TO_NUMBER.get(position)
    return None if n in (None, -1) else n


def apply_position_coloring(prs: Presentation, slide, ordered_positions: List[str], preferred_foot: str) -> None:
    if not ordered_positions:
        return

    main = _resolve_position_number(ordered_positions[0], preferred_foot)
    secondary: List[int] = []
    for p in ordered_positions[1:3]:
        n = _resolve_position_number(p, preferred_foot)
        if n is not None:
            secondary.append(n)

    if main is None and not secondary:
        return

    W = prs.slide_width
    H = prs.slide_height
    x_max = int(W * 0.50)
    y_min = int(H * 0.55)

    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if not (shape.left <= x_max and shape.top >= y_min):
            continue

        txt = (shape.text_frame.text or "").strip()
        if not txt.isdigit():
            continue

        num = int(txt)
        if not (1 <= num <= 11):
            continue

        if main is not None and num == main:
            shape.fill.solid()
            shape.fill.fore_color.rgb = MAIN_BLUE
        elif num in secondary:
            shape.fill.solid()
            shape.fill.fore_color.rgb = SECOND_BLUE


# -----------------------------
# Replacements mapping (includes MV)
# -----------------------------
def build_replacements(
    player: Dict[str, Any],
    transfer_fee: Optional[Dict[str, Any]],
    season_stats: Dict[str, int],
    career_stats: Dict[str, int],
) -> Dict[str, str]:
    info = player.get("info") or {}
    team = player.get("team") or {}
    league = player.get("league") or {}
    contract = player.get("contract") or {}

    name = _as_text(info.get("footballName") or info.get("name") or "")
    dob = _parse_iso_date_to_ddmmyyyy(_as_text(info.get("birthDate")))
    place = _as_text(info.get("birthPlace") or "")

    nats = info.get("nationalities") or []
    nat_names = [str(n.get("name", "")).strip() for n in nats if isinstance(n, dict) and n.get("name")]
    nationalities = ", ".join([n for n in nat_names if n])

    height = _fmt_int(info.get("height"))
    preferred_foot = _as_text(info.get("preferredFoot") or "")
    position = _first_position(info)
    league_name = _as_text(league.get("name") or "")
    club_name = _as_text(team.get("name") or "")

    contract_end = _parse_iso_date_to_ddmmyyyy(_as_text(contract.get("contractEnd")))
    mv = _fmt_money_eur(contract.get("marketValue"))
    tv = _fmt_money_eur(transfer_fee.get("valueEstimateEur")) if transfer_fee else ""

    agency, agent = extract_agent_and_agency(player)

    def variants(token: str) -> List[str]:
        inner = token.strip("{}").strip()
        return list({f"{{{inner}}}", f"{{ {inner} }}", token})

    mapping_pairs = [
        ("{Name}", name),
        ("{ DD/MM/YYYY }", dob),
        ("{Place}", place),
        ("{ Place }", place),
        ("{Nationalities}", nationalities),
        ("{Country}", nationalities),
        ("{ Height }", height),
        ("{ Preferred Foot }", preferred_foot),
        ("{ Preferred foot }", preferred_foot),
        ("{Position}", position),
        ("{ Position }", position),
        ("{League}", league_name),
        ("{ League }", league_name),
        ("{Club}", club_name),
        ("{ Club }", club_name),
        ("{Season_m}", _fmt_int(season_stats.get("matches"))),
        ("{season_min}", _fmt_int(season_stats.get("minutes"))),
        ("{season_g}", _fmt_int(season_stats.get("goals"))),
        ("{season_a}", _fmt_int(season_stats.get("assists"))),
        ("{Career_m}", _fmt_int(career_stats.get("matches"))),
        ("{career_min}", _fmt_int(career_stats.get("minutes"))),
        ("{career_g}", _fmt_int(career_stats.get("goals"))),
        ("{career_a}", _fmt_int(career_stats.get("assists"))),
        ("{con_DD/MM/YYYY}", contract_end),
        ("{TV}", tv),
        ("{MV}", mv),
        ("{Agency}", agency),
        ("{Agent}", agent),
    ]

    out: Dict[str, str] = {}
    for k, v in mapping_pairs:
        for kk in variants(k):
            out[kk] = v
    return out


# -----------------------------
# UI
# -----------------------------
def render_player_card(p: PlayerOption) -> None:
    age_str = "?" if p.age is None else str(p.age)
    st.markdown(
        """
<style>
.player-card {
  border: 1px solid rgba(49, 51, 63, 0.2);
  border-radius: 10px;
  padding: 14px 16px;
  background: rgba(255,255,255,0.02);
}
.player-name { font-size: 16px; font-weight: 700; line-height: 1.2; }
.player-meta { margin-top: 6px; font-size: 13px; opacity: 0.85; }
</style>
""",
        unsafe_allow_html=True,
    )
    st.markdown(
        f"""
<div class="player-card">
  <div class="player-name">{p.name}</div>
  <div class="player-meta">{age_str} - {p.position or "UnknownPos"} - {p.club or "UnknownClub"} ({p.league or "Unknown"})</div>
</div>
""",
        unsafe_allow_html=True,
    )


def main() -> None:
    st.set_page_config(page_title="SciSports Scouting Form Generator", layout="centered")
    st.title("SciSports Scouting Form Generator")

    with st.expander("1) Generate API key (access token)", expanded=True):
        if st.button("Generate API key", type="primary"):
            try:
                token = token_password_grant_from_secrets()
                st.session_state["SCISPORTS_ACCESS_TOKEN"] = token
                os.environ["SCISPORTS_ACCESS_TOKEN"] = token
                st.success("Access token generated and stored for this session.")
            except Exception as e:
                st.error(f"Failed to generate token: {e}")

        if st.session_state.get("SCISPORTS_ACCESS_TOKEN"):
            st.info("Token present ✅ You can now search and select a player.")

    token = st.session_state.get("SCISPORTS_ACCESS_TOKEN", "")
    st.divider()

    st.subheader("2) Search & select player")
    if not token:
        st.warning("Generate the API key first.")
        st.stop()

    with st.form("player_search_form", clear_on_submit=False):
        q = st.text_input("Enter player name", placeholder="e.g. Mees Laros", key="player_search_q")
        submitted = st.form_submit_button("Search")

    if submitted:
        if not q.strip():
            st.warning("Enter a name first.")
        else:
            try:
                _total, options = search_players(token, q)
                st.session_state["player_options"] = options
            except Exception as e:
                st.error(f"Player search failed: {e}")

    options: List[PlayerOption] = st.session_state.get("player_options", [])
    if not options:
        st.info("Search to load players.")
        st.stop()

    st.caption(f"Results: showing {len(options)} (API limit {SEARCH_LIMIT}).")
    selected_label = st.selectbox("Player", options=[o.label() for o in options], index=0)
    selected = next(o for o in options if o.label() == selected_label)
    st.session_state["selected_player_id"] = selected.player_id
    render_player_card(selected)

    st.divider()
    st.subheader("3) Generate Scoutings Form")

    if not os.path.exists(TEMPLATE_PATH):
        st.error(f"Template not found: {TEMPLATE_PATH}. Commit it into the repo root.")
        st.stop()

    if st.button("Generate Scoutings Form", type="primary"):
        with st.spinner("Fetching data and generating PPTX..."):
            try:
                player_id = int(st.session_state["selected_player_id"])
                player = get_player(token, player_id)
                transfer_fee = get_latest_transfer_fee(token, player_id)

                season_label, season_ids, season_stats = compute_target_season_stats(token, player_id, TARGET_SEASON_LABEL)
                career_stats = compute_career_totals(token, player_id)

                if not season_ids:
                    st.warning(f"No season IDs found for season == {TARGET_SEASON_LABEL}. Season stats will be 0.")

                replacements = build_replacements(
                    player=player,
                    transfer_fee=transfer_fee,
                    season_stats=season_stats,
                    career_stats=career_stats,
                )

                positions = (player.get("info") or {}).get("positions") or []
                preferred_foot = _as_text((player.get("info") or {}).get("preferredFoot") or "")

                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    out_path = tmp.name

                fill_pptx(
                    template_path=TEMPLATE_PATH,
                    output_path=out_path,
                    replacements=replacements,
                    positions=[_as_text(p) for p in positions if p],
                    preferred_foot=preferred_foot,
                )

                with open(out_path, "rb") as f:
                    ppt_bytes = f.read()

                safe_name = re.sub(r"[^a-zA-Z0-9_-]+", "_", (player.get("info") or {}).get("name") or "player")[:80]
                out_filename = f"ScoutingsRapport_{safe_name}_{player_id}.pptx"

                st.success("Generated ✅")
                st.caption(f"Season stats used: {season_label} (seasonIds={season_ids})")

                st.download_button(
                    "Download filled Scoutings Rapport (.pptx)",
                    data=ppt_bytes,
                    file_name=out_filename,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                )

                with st.expander("Filled values (debug)"):
                    st.json(replacements)

            except Exception as e:
                st.error(f"Failed to generate PPTX: {e}")


if __name__ == "__main__":
    main()
